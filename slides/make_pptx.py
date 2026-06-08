#!/usr/bin/env python3
"""Generate ML and ANN course PPTX presentations.
Usage:  python slides/make_pptx.py
Output: slides/ml_presentation.pptx
        slides/ann_presentation.pptx
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

OUT = pathlib.Path(__file__).parent

# ── Palette ───────────────────────────────────────────────────────────────
BG    = RGBColor(0x0d,0x11,0x17); TXT  = RGBColor(0xf0,0xf6,0xfc)
GRAY  = RGBColor(0x8b,0x94,0x9e); DGR  = RGBColor(0x48,0x4f,0x58)
BLUE  = RGBColor(0x58,0xa6,0xff); RED  = RGBColor(0xf8,0x51,0x49)
GREEN = RGBColor(0x3f,0xb9,0x50); ORA  = RGBColor(0xd2,0x99,0x22)
BORD  = RGBColor(0x30,0x36,0x3d); PANEL= RGBColor(0x16,0x1b,0x22)
DBLU  = RGBColor(0x0f,0x2a,0x4a); DGRN = RGBColor(0x0f,0x3d,0x20)
DRED  = RGBColor(0x3d,0x0f,0x0f); BLCK = RGBColor(0x0a,0x0e,0x14)

RECT  = 1   # MSO_AUTO_SHAPE_TYPE.RECTANGLE
RRECT = 5   # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE

# ── Helpers ───────────────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def T(slide, text, l, t, w, h, color=TXT, size=18, bold=False,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.color.rgb = color; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    return tb

def R(slide, l, t, w, h, fill=None, lc=BORD, lw=Pt(1.5), st=RRECT):
    sh = slide.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if lc:   sh.line.color.rgb = lc; sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def L(slide, x1, y1, x2, y2, color=BORD, w=Pt(1.2)):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = w
    return c

def hdr(slide, title, course):
    T(slide, course, 0.3, 0.18, 9, 0.32, DGR, 11)
    T(slide, title,  0.5, 0.52, 12.33, 0.85, TXT, 34, bold=True, align=PP_ALIGN.CENTER)
    L(slide, 0, 1.45, 13.33, 1.45)

def box(slide, label, sub, l, t, w=2.0, h=1.0, fc=PANEL, lc=BLUE, lw=Pt(2)):
    R(slide, l, t, w, h, fc, lc, lw)
    T(slide, label, l+0.1, t+0.18, w-0.2, 0.45, TXT, 14, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        T(slide, sub, l+0.1, t+0.58, w-0.2, 0.38, GRAY, 11, align=PP_ALIGN.CENTER)

def arrow(slide, x1, y, x2, color=BLUE):
    L(slide, x1, y, x2, y, color, Pt(2))
    # arrowhead as tiny triangle
    R(slide, x2-0.07, y-0.07, 0.12, 0.14, color, None, st=RECT)

def check_item(slide, text, x, y, color=GREEN):
    R(slide, x, y+0.03, 0.3, 0.3, color, None, lw=Pt(0))
    T(slide, '✓', x+0.06, y+0.05, 0.2, 0.25, BLCK, 14, bold=True, align=PP_ALIGN.CENTER)
    T(slide, text, x+0.42, y, 5.2, 0.36, TXT, 16, bold=True)

def bullet(slide, text, x, y, color=BLUE):
    T(slide, '▶', x, y, 0.28, 0.3, color, 12)
    T(slide, text, x+0.32, y, 5.5, 0.3, GRAY, 15, bold=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  ML PRESENTATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def ml_s1(prs):
    s = blank(prs); set_bg(s)
    # decorative top bar
    R(s, 0, 0, 13.33, 0.06, BLUE, None, st=RECT)
    R(s, 6.67, 0, 6.66, 0.06, RED, None, st=RECT)
    T(s, 'Predictive Maintenance', 0.5, 1.3, 12.33, 1.1,
      TXT, 52, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'via Spiking Neural Networks', 0.5, 2.5, 12.33, 0.8,
      BLUE, 38, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'on Bearing Vibration Data', 0.5, 3.35, 12.33, 0.6,
      GRAY, 22, align=PP_ALIGN.CENTER)
    # waveform hint
    T(s, 'Normal signal  ←——  sparse spikes', 1.0, 4.4, 5.0, 0.4,
      BLUE, 14, align=PP_ALIGN.CENTER)
    T(s, 'Fault signal  ←——  dense spikes', 7.3, 4.4, 5.0, 0.4,
      RED, 14, align=PP_ALIGN.CENTER)
    # sparse spikes visual
    base = [1.3,1.8,2.6,3.3,4.0,4.7,5.5]
    for x in base:
        R(s, x, 4.9, 0.04, 0.35, BLUE, None, lw=Pt(0), st=RECT)
    L(s, 1.0, 5.25, 6.0, 5.25, BORD, Pt(1))
    # dense spikes visual
    for i in range(32):
        R(s, 7.3+i*0.18, 4.9, 0.06, 0.35, RED, None, lw=Pt(0), st=RECT)
    L(s, 7.3, 5.25, 12.5, 5.25, BORD, Pt(1))
    T(s, 'Machine Learning  ·  Final Report  ·  2026',
      0.5, 6.9, 12.33, 0.4, DGR, 14, align=PP_ALIGN.CENTER)

def ml_s2(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Motivation', 'Machine Learning · SNN Predictive Maintenance')
    # 3 problem boxes
    for i,(title,sub,c) in enumerate([
        ('Unplanned\nDowntime',   'No warning before failure', RED),
        ('Manual\nInspection',    'Expensive, error-prone',    ORA),
        ('Early Faults\nMissed',  'Below detection threshold', RED),
    ]):
        x = 0.7 + i*3.6
        R(s, x, 1.6, 3.0, 1.8, DRED if c==RED else RGBColor(0x25,0x1a,0x00), c, Pt(2))
        T(s, title, x+0.15, 1.75, 2.7, 0.8, c, 17, bold=True, align=PP_ALIGN.CENTER)
        T(s, sub,   x+0.15, 2.55, 2.7, 0.6, GRAY, 13, align=PP_ALIGN.CENTER)
    # Down arrow
    T(s, '▼', 6.4, 3.55, 0.6, 0.5, BLUE, 26, align=PP_ALIGN.CENTER)
    # Solution box
    R(s, 1.5, 4.1, 10.33, 1.5, DBLU, BLUE, Pt(2.5))
    T(s, 'Data-Driven SNN Classifier', 1.6, 4.25, 10.1, 0.6,
      BLUE, 26, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'Learn fault signatures from vibration data  ·  <1 KB model  ·  Real-time on microcontroller',
      1.6, 4.85, 10.1, 0.55, GRAY, 15, align=PP_ALIGN.CENTER)
    T(s, 'Constraint: model size < 1 KB  ·  No fixed input window  ·  200 Hz sensor',
      1.6, 5.5, 10.1, 0.4, DGR, 13, align=PP_ALIGN.CENTER)

def ml_s3(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Problem Formulation', 'Machine Learning · SNN Predictive Maintenance')
    # Input
    R(s, 0.4, 1.7, 3.5, 4.2, BLCK, BORD)
    T(s, 'INPUT', 0.4, 1.85, 3.5, 0.35, DGR, 12, align=PP_ALIGN.CENTER)
    T(s, 'Accelerometer\nZ-axis', 0.5, 2.3, 3.3, 0.7, BLUE, 18, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'continuous time-series\nstream (unbounded)', 0.5, 3.1, 3.3, 0.7, GRAY, 13, align=PP_ALIGN.CENTER)
    T(s, '12 kHz · CWRU benchmark', 0.5, 3.85, 3.3, 0.4, DGR, 12, align=PP_ALIGN.CENTER)
    # Arrow + Model label
    T(s, '──────►', 4.1, 3.65, 1.2, 0.4, BLUE, 18, align=PP_ALIGN.CENTER)
    T(s, 'Model', 4.1, 3.3, 1.2, 0.35, GRAY, 12, align=PP_ALIGN.CENTER)
    # Output
    R(s, 5.5, 2.1, 3.5, 1.4, DGRN, GREEN, Pt(2))
    T(s, '✓  Normal', 5.6, 2.4, 3.3, 0.6, GREEN, 22, bold=True, align=PP_ALIGN.CENTER)
    R(s, 5.5, 3.85, 3.5, 1.4, DRED, RED, Pt(2))
    T(s, '✗  Fault', 5.6, 4.15, 3.3, 0.6, RED, 22, bold=True, align=PP_ALIGN.CENTER)
    # Constraints
    R(s, 9.5, 1.7, 3.4, 4.2, PANEL, BORD)
    T(s, 'Constraints', 9.6, 1.85, 3.2, 0.35, ORA, 14, bold=True)
    for i,c in enumerate([
        '• Unbounded input (streaming)',
        '• Model size  < 1 KB',
        '• Real-time inference',
        '• No fixed window size',
        '• Binary output per interval',
    ]):
        T(s, c, 9.7, 2.35+i*0.5, 3.1, 0.45, GRAY, 13)

def ml_s4(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'CWRU Bearing Dataset', 'Machine Learning · SNN Predictive Maintenance')
    R(s, 0.4, 1.55, 6.2, 0.5, DBLU, BLUE, Pt(1))
    T(s, 'Case Western Reserve University  ·  Public Benchmark',
      0.5, 1.63, 6.0, 0.35, BLUE, 13, bold=True, align=PP_ALIGN.CENTER)
    # File cards
    T(s, 'NORMAL', 0.5, 2.2, 2.5, 0.3, GRAY, 12, bold=True)
    R(s, 0.4, 2.5, 5.8, 0.75, DGRN, GREEN, Pt(1.5))
    T(s, '97.mat   —   Normal baseline   ·   243,938 samples', 0.55, 2.65, 5.5, 0.4, GREEN, 14, bold=True)
    T(s, 'INNER RACE FAULT', 0.5, 3.4, 3.0, 0.3, GRAY, 12, bold=True)
    for i,(f,d) in enumerate([('105.mat','0.007 in.'),('118.mat','0.014 in.'),('130.mat','0.021 in.')]):
        R(s, 0.4, 3.7+i*0.82, 5.8, 0.7, DRED, RED, Pt(1.5))
        T(s, f'{f}   —   fault diameter {d}', 0.55, 3.85+i*0.82, 5.5, 0.4, RED, 14, bold=True)
    # Bar chart (manual)
    T(s, 'Sample Count', 7.5, 1.6, 5.5, 0.35, GRAY, 14, bold=True, align=PP_ALIGN.CENTER)
    base_x, base_y = 7.2, 6.5
    # Normal bar: 243938 → full width ~5"
    R(s, base_x+0.3, base_y-1.22, 5.0, 1.22, DGRN, GREEN, Pt(1), st=RECT)
    T(s, '243,938', base_x+5.4, base_y-1.0, 1.5, 0.35, GREEN, 14, bold=True)
    T(s, 'Normal', base_x+0.3, base_y+0.05, 2.0, 0.3, GREEN, 12, bold=True)
    # Fault bar: 121265 → ~2.49"
    R(s, base_x+0.3, base_y-0.7, 2.5, 0.7, DRED, RED, Pt(1), st=RECT)
    T(s, '121,265', base_x+2.9, base_y-0.5, 1.5, 0.35, RED, 14, bold=True)
    T(s, 'Fault', base_x+0.3, base_y+0.05+0.35, 2.0, 0.3, RED, 12, bold=True)
    L(s, base_x+0.3, base_y, base_x+5.5, base_y, BORD, Pt(1.5))
    T(s, 'Sampling Rate: 12,000 Hz  ·  Drive-End Accelerometer  ·  scipy.io.loadmat()',
      0.4, 7.0, 12.5, 0.35, DGR, 12, align=PP_ALIGN.CENTER)

def ml_s5(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Delta Modulation → Spike Encoding', 'Machine Learning · SNN Predictive Maintenance')
    T(s, 'spike(t) = 1  if  |accel(t) − accel(t−1)| > θ        θ = 0.05',
      1.5, 1.5, 10.33, 0.45, GREEN, 17, bold=True, italic=True, align=PP_ALIGN.CENTER)
    # Normal column
    T(s, 'NORMAL  Bearing', 0.5, 2.05, 5.8, 0.38, BLUE, 15, bold=True, align=PP_ALIGN.CENTER)
    R(s, 0.4, 2.45, 5.8, 1.3, BLCK, BORD)
    T(s, 'Raw signal  (small amplitude oscillation)', 0.55, 2.55, 5.5, 0.35, DGR, 11)
    # Wavy line hint (alternating tiny rects)
    for i in range(18):
        h = 0.18 if i%2==0 else 0.04
        R(s, 0.6+i*0.3, 3.3-h/2, 0.22, h, BLUE, None, lw=Pt(0), st=RECT)
    R(s, 0.4, 3.85, 5.8, 1.0, BLCK, BORD)
    T(s, 'Spike train  (sparse)', 0.55, 3.92, 5.5, 0.3, DGR, 11)
    L(s, 0.5, 4.65, 6.1, 4.65, RGBColor(0x21,0x26,0x2d), Pt(1))
    for x in [0.9, 1.7, 2.6, 3.5, 4.1, 4.9, 5.7]:
        R(s, x, 4.35, 0.055, 0.3, BLUE, None, lw=Pt(0), st=RECT)
    R(s, 0.9, 5.0, 4.8, 0.65, DBLU, BLUE, Pt(2))
    T(s, 'Spike Rate', 1.0, 5.08, 4.5, 0.28, GRAY, 13, align=PP_ALIGN.CENTER)
    T(s, '0.185', 1.0, 5.35, 4.5, 0.25, BLUE, 22, bold=True, align=PP_ALIGN.CENTER)
    # Fault column
    T(s, 'FAULT  Bearing  (Inner Race)', 7.1, 2.05, 5.8, 0.38, RED, 15, bold=True, align=PP_ALIGN.CENTER)
    R(s, 7.0, 2.45, 5.8, 1.3, BLCK, BORD)
    T(s, 'Raw signal  (large amplitude impulses)', 7.15, 2.55, 5.5, 0.35, DGR, 11)
    for i in range(12):
        h = 0.55 if i%2==0 else 0.08
        R(s, 7.1+i*0.45, 3.2-h/2, 0.3, h, RED, None, lw=Pt(0), st=RECT)
    R(s, 7.0, 3.85, 5.8, 1.0, BLCK, BORD)
    T(s, 'Spike train  (dense)', 7.15, 3.92, 5.5, 0.3, DGR, 11)
    L(s, 7.1, 4.65, 12.7, 4.65, RGBColor(0x21,0x26,0x2d), Pt(1))
    for i in range(58):
        R(s, 7.15+i*0.095, 4.35, 0.055, 0.3, RED, None, lw=Pt(0), st=RECT)
    R(s, 7.6, 5.0, 4.8, 0.65, DRED, RED, Pt(2))
    T(s, 'Spike Rate', 7.7, 5.08, 4.5, 0.28, GRAY, 13, align=PP_ALIGN.CENTER)
    T(s, '0.854', 7.7, 5.35, 4.5, 0.25, RED, 22, bold=True, align=PP_ALIGN.CENTER)
    T(s, '4.6×  difference  →  directly classifiable by spike rate',
      1.5, 5.85, 10.33, 0.42, GREEN, 18, bold=True, align=PP_ALIGN.CENTER)
    # VS divider
    L(s, 6.67, 2.0, 6.67, 5.8, BORD, Pt(1))

def ml_s6(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Model Architecture — Streaming SNN', 'Machine Learning · SNN Predictive Maintenance')
    labels = [
        ('Spike\nInput',    '1 scalar\nper step', BLUE,  BLCK),
        ('FC1\nLinear',     '1 → 32\nno bias',   BLUE,  DBLU),
        ('LIF1\nLeaky',     'β learned\npersistent mem', GREEN, DGRN),
        ('FC2\nLinear',     '32 → 2\nno bias',   BLUE,  DBLU),
        ('Σ Accum.\n200 steps', '≈1 second\nrate coding', ORA, PANEL),
        ('argmax\nClassify','Normal\nor Fault',  GREEN, DGRN),
    ]
    x0, y0, bw, bh, gap = 0.35, 2.2, 1.9, 1.8, 0.32
    for i,(lbl,sub,lc,fc) in enumerate(labels):
        bx = x0 + i*(bw+gap)
        R(s, bx, y0, bw, bh, fc, lc, Pt(2.5))
        T(s, lbl, bx+0.1, y0+0.25, bw-0.2, 0.65, TXT, 15, bold=True, align=PP_ALIGN.CENTER)
        T(s, sub, bx+0.1, y0+0.95, bw-0.2, 0.65, GRAY, 11, align=PP_ALIGN.CENTER)
        if i < len(labels)-1:
            ax = bx+bw+0.02
            L(s, ax, y0+bh/2, ax+gap-0.04, y0+bh/2, BLUE, Pt(2))
    # LIF details
    R(s, 2.87+2*(1.9+0.32), 4.2, 1.9, 1.0, DGRN, GREEN, Pt(1.5))
    T(s, 'mem = β·mem + W·x\nfire if mem ≥ θ;\nmem −= θ', 2.87+2*(1.9+0.32)+0.1, 4.28, 1.7, 0.85, GREEN, 10, italic=True)
    # Param badge
    R(s, 3.5, 4.3, 6.33, 0.75, DGRN, GREEN, Pt(2.5))
    T(s, 'Total: 97 parameters  ·  0.38 KB (float32)',
      3.6, 4.5, 6.1, 0.42, GREEN, 18, bold=True, align=PP_ALIGN.CENTER)
    # Persistent mem note
    T(s, '⟵  mem₁ never reset between samples  —  persistent temporal memory  ⟶',
      1.5, 5.25, 10.33, 0.4, GREEN, 14, italic=True, align=PP_ALIGN.CENTER)

def ml_s7(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Leaky Integrate-and-Fire Neuron', 'Machine Learning · SNN Predictive Maintenance')
    # Equation
    R(s, 0.4, 1.6, 7.5, 1.0, PANEL, GREEN, Pt(2))
    T(s, 'mem(t) = β · mem(t−1)  +  W · spike(t)', 0.55, 1.72, 7.2, 0.42, GREEN, 17, bold=True, italic=True)
    T(s, 'if mem(t) ≥ θ :   spike!   mem(t) −= θ   (soft reset)', 0.55, 2.12, 7.2, 0.38, TXT, 15, italic=True)
    # Parameters table
    R(s, 0.4, 2.75, 7.5, 3.5, PANEL, BORD)
    T(s, 'Parameters', 0.55, 2.88, 7.2, 0.38, ORA, 14, bold=True)
    L(s, 0.4, 3.3, 7.9, 3.3, BORD, Pt(1))
    rows = [
        ('β  (beta)',      'leak factor ∈ (0,1)',  'controls decay speed  —  LEARNED',  BLUE),
        ('θ  (threshold)', 'fire threshold = 1.0',  'fixed; soft reset: mem −= θ (not zeroed)', ORA),
        ('W  (weights)',   'synaptic weights',      'learned via surrogate BPTT, no bias terms', GREEN),
        ('x(t)',           'binary spike {0, 1}',   'from delta modulation of accelerometer',    GRAY),
    ]
    for i,(n,v,d,c) in enumerate(rows):
        y = 3.4 + i*0.65
        T(s, n, 0.6, y, 1.6, 0.55, c, 13, bold=True)
        T(s, v, 2.3, y, 2.2, 0.55, TXT, 13)
        T(s, d, 4.6, y, 3.6, 0.55, GRAY, 12)
    # Membrane potential diagram (simplified line chart)
    try:
        cd = ChartData()
        cd.categories = list(range(1,11))
        cd.add_series('mem(t)', (0.05,0.18,0.38,0.64,0.88,1.05,0.05,0.22,0.50,0.82))
        cf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(8.2), Inches(1.6), Inches(4.9), Inches(4.6), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        ch.plots[0].series[0].format.line.color.rgb = GREEN
        ch.plots[0].series[0].format.line.width = Pt(2.5)
        ch.chart_area.format.fill.solid()
        ch.chart_area.format.fill.fore_color.rgb = BLCK
        ch.plot_area.format.fill.solid()
        ch.plot_area.format.fill.fore_color.rgb = BLCK
    except Exception:
        T(s, '↗ membrane rises → fires at θ=1.0 → soft reset → rises again',
          8.2, 3.5, 4.9, 0.5, GREEN, 13, align=PP_ALIGN.CENTER)
    T(s, 'threshold θ = 1.0  (dashed line on chart)', 8.2, 6.3, 4.9, 0.4, ORA, 13, italic=True, align=PP_ALIGN.CENTER)

def ml_s8(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Training Pipeline', 'Machine Learning · SNN Predictive Maintenance')
    items = [
        ('CWRU\n.mat', '12 kHz\nraw', BLUE),
        ('Delta\nMod', 'θ=0.05\nbinary', BLUE),
        ('Chunks\nT=256', 'TBPTT\nboundary', ORA),
        ('Streaming\nSNN', '97 params\nFC+LIF+FC', GREEN),
        ('Cross-\nEntropy', 'surrogate\ngrad', BLUE),
        ('Adam', 'lr=1e−3\noptimizer', BLUE),
        ('Early\nStop', 'patience\n=10', GREEN),
    ]
    x0, y0, bw, bh, gap = 0.32, 2.1, 1.7, 1.65, 0.2
    for i,(lbl,sub,lc) in enumerate(items):
        bx = x0 + i*(bw+gap)
        R(s, bx, y0, bw, bh, DBLU if lc==BLUE else (DGRN if lc==GREEN else PANEL), lc, Pt(2))
        T(s, lbl, bx+0.1, y0+0.2, bw-0.2, 0.65, TXT, 14, bold=True, align=PP_ALIGN.CENTER)
        T(s, sub, bx+0.1, y0+0.9, bw-0.2, 0.6, GRAY, 11, align=PP_ALIGN.CENTER)
        if i < len(items)-1:
            ax = bx+bw+0.02
            L(s, ax, y0+bh/2, ax+gap-0.04, y0+bh/2, BLUE, Pt(2))
    # TBPTT explanation
    R(s, 0.4, 4.0, 7.5, 1.65, PANEL, ORA, Pt(1.5))
    T(s, 'Truncated BPTT', 0.55, 4.1, 7.2, 0.38, ORA, 14, bold=True)
    for i,t in enumerate(['Unroll C=256 steps per chunk',
                          'Detach gradient at chunk boundary  (stop grad)',
                          'mem₁ persists across all chunks during INFERENCE']):
        T(s, f'• {t}', 0.7, 4.52+i*0.38, 7.0, 0.35, GRAY, 13)
    # Train/test split
    R(s, 8.2, 4.0, 4.8, 1.65, PANEL, BORD)
    T(s, 'Data Split', 8.35, 4.1, 4.5, 0.38, GRAY, 14, bold=True)
    T(s, '80% train  →  1,140 chunks', 8.35, 4.52, 4.5, 0.38, BLUE, 14)
    T(s, '20% test   →     285 chunks', 8.35, 4.9, 4.5, 0.38, GREEN, 14)

def ml_s9(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Training Results', 'Machine Learning · SNN Predictive Maintenance')
    # Line chart for loss
    try:
        cd = ChartData()
        cd.categories = ['Ep 1','Ep 2','Ep 3','Ep 4','Ep 5']
        cd.add_series('Loss', (1.4727, 0.50, 0.09, 0.01, 0.0001))
        cf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.3), cd)
        ch = cf.chart; ch.has_title = True; ch.has_legend = False
        ch.chart_title.text_frame.text = 'Training Loss'
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = GRAY
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        ch.plots[0].series[0].format.line.color.rgb = GREEN
        ch.plots[0].series[0].format.line.width = Pt(2.5)
        ch.chart_area.format.fill.solid(); ch.chart_area.format.fill.fore_color.rgb = BLCK
        ch.plot_area.format.fill.solid(); ch.plot_area.format.fill.fore_color.rgb = BLCK
    except Exception:
        R(s, 0.5, 1.6, 6.0, 5.3, BLCK, BORD)
        T(s, 'Loss: 1.47 → 0.50 → 0.09 → 0.01 → 0.0001\n(converges in 5 epochs)',
          0.8, 3.5, 5.4, 1.0, GREEN, 16, align=PP_ALIGN.CENTER)
    # Accuracy badge
    R(s, 7.0, 2.0, 5.9, 4.5, DGRN, GREEN, Pt(3))
    T(s, '100%', 7.1, 2.3, 5.7, 1.8, GREEN, 80, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'Test Accuracy', 7.1, 4.1, 5.7, 0.55, TXT, 22, bold=True, align=PP_ALIGN.CENTER)
    T(s, '285 test chunks', 7.1, 4.65, 5.7, 0.4, GRAY, 16, align=PP_ALIGN.CENTER)
    # Classification table
    R(s, 0.5, 6.55, 12.33, 0.7, PANEL, BORD)
    T(s, 'Fault:   P=1.00   R=1.00   F1=1.00  (102 chunks)       Normal:   P=1.00   R=1.00   F1=1.00  (183 chunks)',
      0.7, 6.68, 12.0, 0.42, GREEN, 14, bold=True, align=PP_ALIGN.CENTER)

def ml_s10(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Why Does It Work?', 'Machine Learning · SNN Predictive Maintenance')
    # Spike rate bar chart
    T(s, 'Spike Rate per Sample', 0.5, 1.55, 6.0, 0.38, GRAY, 14, bold=True, align=PP_ALIGN.CENTER)
    bx, by, bw = 0.9, 1.98, 1.8
    chart_h = 4.5
    # Normal bar: 0.185 * 4.5 = 0.83"
    nh = 0.185 * chart_h
    R(s, bx, by+(chart_h-nh), bw, nh, DBLU, BLUE, Pt(1.5), st=RECT)
    T(s, '0.185', bx, by+(chart_h-nh)-0.35, bw, 0.32, BLUE, 18, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'Normal', bx, by+chart_h+0.08, bw, 0.32, BLUE, 14, bold=True, align=PP_ALIGN.CENTER)
    # Fault bar: 0.854 * 4.5 = 3.84"
    fh = 0.854 * chart_h
    R(s, bx+2.5, by+(chart_h-fh), bw, fh, DRED, RED, Pt(1.5), st=RECT)
    T(s, '0.854', bx+2.5, by+(chart_h-fh)-0.35, bw, 0.32, RED, 18, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'Fault', bx+2.5, by+chart_h+0.08, bw, 0.32, RED, 14, bold=True, align=PP_ALIGN.CENTER)
    L(s, 0.5, by+chart_h, 6.0, by+chart_h, BORD, Pt(1.5))
    T(s, '4.6×', 2.15, by+1.8, 1.2, 0.5, GREEN, 32, bold=True, align=PP_ALIGN.CENTER)
    # Explanation cards
    for i,(title,body,c) in enumerate([
        ('① Delta Modulation',
         'Fault impulses → |Δaccel| exceeds θ nearly every step → dense spikes', BLUE),
        ('② LIF as Temporal Integrator',
         'mem(t) accumulates spike rate → robust, persistent memory, noise-resistant', GREEN),
        ('③ Rate Coding Output',
         'Logits summed over 200 steps → argmax → robust to single-step noise', ORA),
    ]):
        R(s, 6.5, 1.6+i*1.82, 6.5, 1.65, PANEL, c, Pt(1.5))
        T(s, title, 6.65, 1.72+i*1.82, 6.2, 0.38, c, 14, bold=True)
        T(s, body,  6.65, 2.12+i*1.82, 6.2, 0.7,  GRAY, 13)

def ml_s11(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Model Efficiency', 'Machine Learning · SNN Predictive Maintenance')
    T(s, 'Parameter Count Comparison', 0.5, 1.55, 12.33, 0.38, GRAY, 15, bold=True, align=PP_ALIGN.CENTER)
    # Horizontal bars
    x0, scale = 2.8, 10.2/10000  # 10.2" for 10000 params
    data = [
        ('1D-CNN',      10000, ORA,   DRED),
        ('LSTM',         5000, GRAY,  PANEL),
        ('SNN  (ours)',    97, GREEN, DGRN),
    ]
    for i,(name,params,lc,fc) in enumerate(data):
        y = 2.15 + i*1.55
        T(s, name, 0.5, y+0.3, 2.2, 0.45, lc, 15, bold=True)
        bw = max(params*scale, 0.12)
        R(s, x0, y+0.1, bw, 0.75, fc, lc, Pt(1.5), st=RECT)
        T(s, f'{params:,}', x0+bw+0.1, y+0.28, 1.8, 0.4, lc, 16, bold=True)
    # SNN callout
    R(s, 8.5, 5.5, 4.3, 1.5, DGRN, GREEN, Pt(3))
    T(s, '97', 8.6, 5.6, 4.1, 0.7, GREEN, 60, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'parameters  ·  0.38 KB', 8.6, 6.3, 4.1, 0.45, GRAY, 16, align=PP_ALIGN.CENTER)
    T(s, '100× fewer parameters than a typical CNN  ·  no multiplications for binary input',
      0.5, 6.85, 12.33, 0.42, TXT, 16, bold=True, align=PP_ALIGN.CENTER)

def ml_s12(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Conclusion', 'Machine Learning · SNN Predictive Maintenance')
    T(s, 'Achieved', 0.5, 1.55, 5.8, 0.38, GREEN, 15, bold=True)
    for i,t in enumerate([
        '100% test accuracy on CWRU benchmark',
        'Only 97 parameters  ·  0.38 KB model',
        'True streaming — no fixed window size',
        'Delta modulation: effective spike-rate feature',
        'Persistent LIF memory = temporal integration',
    ]):
        check_item(s, t, 0.5, 2.0+i*0.7, GREEN)
    L(s, 6.3, 1.55, 6.3, 7.1, BORD, Pt(1.5))
    T(s, 'Future Work', 6.8, 1.55, 6.1, 0.38, BLUE, 15, bold=True)
    for i,(t,d) in enumerate([
        ('Multi-class fault detection', 'ball fault · outer race fault · combined'),
        ('Real motor data at 200 Hz',   'validate vs CWRU 12 kHz performance'),
        ('INT8 quantization',           '4× compression · faster MCU inference'),
        ('Multi-axis + current sensor', 'richer vibration feature space'),
    ]):
        bullet(s, t, 6.8, 2.0+i*1.12)
        T(s, d, 7.15, 2.32+i*1.12, 5.8, 0.3, DGR, 13)
    R(s, 0.4, 6.65, 12.53, 0.7, DGRN, GREEN, Pt(2.5))
    T(s, '97-parameter streaming SNN achieves fault detection on par with much larger models',
      0.6, 6.77, 12.1, 0.42, GREEN, 15, bold=True, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  ANN PRESENTATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def ann_s1(prs):
    s = blank(prs); set_bg(s)
    R(s, 0, 0, 13.33, 0.06, GREEN, None, st=RECT)
    R(s, 6.67, 0, 6.66, 0.06, BLUE, None, st=RECT)
    T(s, 'Streaming Spiking Neural Networks', 0.5, 1.2, 12.33, 1.1,
      TXT, 48, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'for Real-Time Fault Detection', 0.5, 2.4, 12.33, 0.8,
      GREEN, 36, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'LIF Neurons  ·  Surrogate Gradient  ·  TBPTT  ·  Edge Deployment',
      0.5, 3.3, 12.33, 0.55, GRAY, 20, align=PP_ALIGN.CENTER)
    T(s, 'Artificial Neural Networks  ·  Final Report  ·  2026',
      0.5, 4.1, 12.33, 0.45, DGR, 16, align=PP_ALIGN.CENTER)
    T(s, 'mem(t) = β · mem(t−1) + W · spike(t)    ·    fire if mem ≥ θ    ·    97 parameters',
      0.5, 6.9, 12.33, 0.42, RGBColor(0x21,0x26,0x2d), 14, italic=True, align=PP_ALIGN.CENTER)

def ann_s2(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'From ANN to Spiking Neural Network', 'ANN Course · SNN Predictive Maintenance')
    # Table
    tbl = s.shapes.add_table(6, 3, Inches(0.4), Inches(1.6), Inches(12.53), Inches(5.65)).table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(5.0)
    tbl.columns[2].width = Inches(5.03)
    headers = ['Property', 'Traditional ANN', 'Spiking Neural Network (SNN)']
    hcolors = [BORD, ORA, GREEN]
    rows_data = [
        ('Activation',       'Continuous  (0.0 – 1.0)',          'Binary spikes  (0 or 1)'),
        ('Communication',    'Dense float vectors  (all active)', 'Sparse events  (only firing)'),
        ('Temporal State',   'Stateless  (no memory)',            'Membrane potential  (persistent)'),
        ('Gradient Method',  'Standard backprop  (differentiable)','Surrogate gradient  (non-diff. spike)'),
        ('Edge Efficiency',  'Heavy  (float multiply per step)',  'Light  (binary: add only, no multiply)'),
    ]
    for col,(hdr_txt, hc) in enumerate(zip(headers, hcolors)):
        cell = tbl.cell(0, col)
        cell.text = hdr_txt
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = hc
        p.runs[0].font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    for row,(prop,ann_v,snn_v) in enumerate(rows_data):
        fc = BLCK if row%2==0 else PANEL
        for col,(val,vc) in enumerate([(prop,GRAY),(ann_v,TXT),(snn_v,GREEN)]):
            cell = tbl.cell(row+1, col)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = fc
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.color.rgb = vc
            p.runs[0].font.size = Pt(13)
            p.runs[0].font.bold = (col==2)
            p.alignment = PP_ALIGN.CENTER if col>0 else PP_ALIGN.LEFT

def ann_s3(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Leaky Integrate-and-Fire Neuron', 'ANN Course · SNN Predictive Maintenance')
    R(s, 0.4, 1.6, 8.0, 0.85, PANEL, GREEN, Pt(2))
    T(s, 'mem(t) = β · mem(t−1)  +  W · spike(t)', 0.55, 1.72, 7.7, 0.42, GREEN, 17, bold=True, italic=True)
    T(s, 'if mem(t) ≥ θ :   emit spike ;   mem(t) −= θ   ← soft reset (not zero)', 0.55, 2.1, 7.7, 0.38, TXT, 14, italic=True)
    # Parameters
    R(s, 0.4, 2.6, 8.0, 3.8, PANEL, BORD)
    T(s, 'Parameter', 0.55, 2.72, 2.0, 0.38, ORA, 13, bold=True)
    T(s, 'Value / Range', 2.6, 2.72, 2.2, 0.38, ORA, 13, bold=True)
    T(s, 'Role', 4.9, 2.72, 3.4, 0.38, ORA, 13, bold=True)
    L(s, 0.4, 3.15, 8.4, 3.15, BORD, Pt(1))
    params = [
        ('β  (beta)', '(0, 1) — LEARNED', 'leak factor · controls decay speed'),
        ('θ  (threshold)', '1.0 — fixed', 'fire level · soft reset preserves energy'),
        ('W  (weights)', '[32,1] matrix', 'synaptic · learned via surrogate BPTT'),
        ('x(t)', '{0, 1} binary', 'delta modulation of accelerometer'),
    ]
    for i,(n,v,d) in enumerate(params):
        y = 3.25+i*0.72
        T(s, n, 0.6, y, 2.0, 0.55, BLUE, 13, bold=True)
        T(s, v, 2.6, y, 2.2, 0.55, TXT, 13)
        T(s, d, 4.9, y, 3.4, 0.55, GRAY, 12)
    # Membrane potential chart
    try:
        cd = ChartData()
        cd.categories = [f't{i}' for i in range(12)]
        cd.add_series('mem', (0.02,0.16,0.35,0.58,0.80,1.05,0.05,0.18,0.40,0.68,0.92,0.08))
        cf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(8.7), Inches(1.6), Inches(4.3), Inches(4.8), cd)
        ch = cf.chart; ch.has_title = True; ch.has_legend = False
        ch.chart_title.text_frame.text = 'mem(t) dynamics'
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = GRAY
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        ch.plots[0].series[0].format.line.color.rgb = GREEN
        ch.plots[0].series[0].format.line.width = Pt(2.5)
        ch.chart_area.format.fill.solid(); ch.chart_area.format.fill.fore_color.rgb = BLCK
        ch.plot_area.format.fill.solid(); ch.plot_area.format.fill.fore_color.rgb = BLCK
    except Exception:
        R(s, 8.7, 1.6, 4.3, 4.8, BLCK, BORD)
        T(s, 'mem rises → fires at θ=1.0\nsoft reset → rises again',
          8.9, 3.5, 3.9, 1.0, GREEN, 15, align=PP_ALIGN.CENTER)
    T(s, 'Soft reset (mem−=θ) preserves sub-threshold energy — more info than hard zero',
      0.4, 6.55, 12.53, 0.42, DGR, 13, italic=True, align=PP_ALIGN.CENTER)

def ann_s4(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Delta Modulation — Spike Encoding', 'ANN Course · SNN Predictive Maintenance')
    R(s, 0.4, 1.55, 5.8, 0.65, PANEL, BLUE, Pt(1.5))
    T(s, 'Biological Inspiration: retinal ganglion cells fire on luminance CHANGE, not absolute level',
      0.55, 1.65, 5.55, 0.45, BLUE, 12, italic=True)
    R(s, 6.5, 1.55, 6.4, 0.65, DGRN, GREEN, Pt(2))
    T(s, 'spike(t) = 1  if  |Δaccel| > θ = 0.05', 6.65, 1.65, 6.1, 0.45, GREEN, 14, bold=True, italic=True)
    # Same visual as ML slide 5 (reuse layout)
    ml_s5_visual_only(s)

def ml_s5_visual_only(s):
    # Reuse the spike visual from ml_s5 without redoing the header
    T(s, 'NORMAL  Bearing', 0.5, 2.35, 5.8, 0.38, BLUE, 14, bold=True, align=PP_ALIGN.CENTER)
    R(s, 0.4, 2.75, 5.8, 1.1, BLCK, BORD)
    for i in range(18):
        h = 0.18 if i%2==0 else 0.04
        R(s, 0.6+i*0.3, 3.55-h/2, 0.22, h, BLUE, None, lw=Pt(0), st=RECT)
    R(s, 0.4, 3.95, 5.8, 0.9, BLCK, BORD)
    L(s, 0.5, 4.65, 6.1, 4.65, RGBColor(0x21,0x26,0x2d), Pt(1))
    for x in [0.9,1.7,2.6,3.5,4.1,4.9,5.7]:
        R(s, x, 4.18, 0.055, 0.3, BLUE, None, lw=Pt(0), st=RECT)
    R(s, 0.9, 5.0, 4.8, 0.65, DBLU, BLUE, Pt(2))
    T(s, 'Spike Rate  0.185', 1.0, 5.15, 4.5, 0.4, BLUE, 20, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'FAULT  Bearing', 7.1, 2.35, 5.8, 0.38, RED, 14, bold=True, align=PP_ALIGN.CENTER)
    R(s, 7.0, 2.75, 5.8, 1.1, BLCK, BORD)
    for i in range(12):
        h = 0.5 if i%2==0 else 0.08
        R(s, 7.1+i*0.45, 3.5-h/2, 0.3, h, RED, None, lw=Pt(0), st=RECT)
    R(s, 7.0, 3.95, 5.8, 0.9, BLCK, BORD)
    L(s, 7.1, 4.65, 12.7, 4.65, RGBColor(0x21,0x26,0x2d), Pt(1))
    for i in range(55):
        R(s, 7.15+i*0.1, 4.18, 0.065, 0.3, RED, None, lw=Pt(0), st=RECT)
    R(s, 7.6, 5.0, 4.8, 0.65, DRED, RED, Pt(2))
    T(s, 'Spike Rate  0.854', 7.7, 5.15, 4.5, 0.4, RED, 20, bold=True, align=PP_ALIGN.CENTER)
    T(s, '4.6× difference  →  directly classifiable',
      1.5, 5.9, 10.33, 0.42, GREEN, 17, bold=True, align=PP_ALIGN.CENTER)
    L(s, 6.67, 2.3, 6.67, 5.7, BORD, Pt(1))

def ann_s5(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Windowed SNN  vs  Streaming SNN', 'ANN Course · SNN Predictive Maintenance')
    # Left: Windowed (bad)
    R(s, 0.3, 1.55, 5.9, 0.55, DRED, RED, Pt(2))
    T(s, '❌  Windowed  (old approach)', 0.5, 1.65, 5.5, 0.38, RED, 15, bold=True, align=PP_ALIGN.CENTER)
    R(s, 0.3, 2.2, 5.9, 1.0, BLCK, BORD)
    T(s, 'Buffer: [ x₁, x₂, ..., x₁₂₈ ]', 0.5, 2.35, 5.5, 0.38, GRAY, 14, align=PP_ALIGN.CENTER)
    T(s, 'must wait 128 samples before first output', 0.5, 2.7, 5.5, 0.35, DGR, 12, align=PP_ALIGN.CENTER)
    R(s, 0.8, 3.35, 4.9, 0.85, DRED, RED, Pt(1.5))
    T(s, 'FC1:  Linear(128 → 32)', 1.0, 3.5, 4.5, 0.38, RED, 15, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'WINDOW size baked into weight matrix!', 1.0, 3.75, 4.5, 0.35, GRAY, 12, align=PP_ALIGN.CENTER)
    R(s, 0.3, 4.35, 5.9, 1.6, PANEL, RED, Pt(1))
    T(s, 'Problems:', 0.5, 4.48, 5.5, 0.35, RED, 13, bold=True)
    for i,t in enumerate(['• Changing window = full retrain',
                          '• Waits 128 samples before output',
                          '• Window boundary artifacts']):
        T(s, t, 0.55, 4.85+i*0.35, 5.5, 0.32, GRAY, 12)
    R(s, 0.8, 6.1, 4.9, 0.65, PANEL, BORD)
    T(s, 'Params:  128×32 + 32×2 = 4,160  ·  16 KB', 1.0, 6.25, 4.5, 0.38, RED, 13, align=PP_ALIGN.CENTER)
    L(s, 6.5, 1.55, 6.5, 7.1, BORD, Pt(2))
    T(s, 'VS', 6.15, 4.1, 0.7, 0.5, GRAY, 20, bold=True, align=PP_ALIGN.CENTER)
    # Right: Streaming (good)
    R(s, 6.8, 1.55, 6.1, 0.55, DGRN, GREEN, Pt(2.5))
    T(s, '✓  Streaming  (this work)', 7.0, 1.65, 5.7, 0.38, GREEN, 15, bold=True, align=PP_ALIGN.CENTER)
    R(s, 6.8, 2.2, 6.1, 1.0, BLCK, GREEN, Pt(1))
    T(s, 'one spike scalar per timestep', 7.0, 2.38, 5.7, 0.38, GREEN, 14, align=PP_ALIGN.CENTER)
    T(s, 'instant output, no buffer needed', 7.0, 2.72, 5.7, 0.35, DGR, 12, align=PP_ALIGN.CENTER)
    R(s, 7.1, 3.35, 5.3, 0.85, DGRN, GREEN, Pt(2))
    T(s, 'FC1:  Linear(1 → 32)', 7.3, 3.5, 4.9, 0.38, GREEN, 15, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'no window dependency, scalar input', 7.3, 3.75, 4.9, 0.35, GRAY, 12, align=PP_ALIGN.CENTER)
    R(s, 6.8, 4.35, 6.1, 1.6, DGRN, GREEN, Pt(1))
    T(s, 'Key Innovation:', 7.0, 4.48, 5.7, 0.35, GREEN, 13, bold=True)
    for i,t in enumerate(['• mem₁ never reset between samples',
                          '• Temporal context spans entire stream',
                          '• Decision every 200 steps (≈1 second)']):
        T(s, t, 7.0, 4.85+i*0.35, 5.7, 0.32, GRAY, 12)
    R(s, 7.1, 6.1, 5.3, 0.65, DGRN, GREEN, Pt(2))
    T(s, 'Params:  32 + 64 + 1 = 97  ·  0.38 KB', 7.3, 6.25, 4.9, 0.38, GREEN, 13, bold=True, align=PP_ALIGN.CENTER)
    T(s, '43× fewer parameters · truly real-time · no window coupling',
      0.5, 7.05, 12.33, 0.35, TXT, 15, bold=True, align=PP_ALIGN.CENTER)

def ann_s6(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Network Architecture', 'ANN Course · SNN Predictive Maintenance')
    R(s, 1.0, 1.55, 11.33, 0.52, PANEL, GREEN, Pt(1.5))
    T(s, 'cur₁ = FC1(spike)  →  spk₁, mem₁ = LIF1(cur₁, mem₁)  →  logit += FC2(spk₁)',
      1.1, 1.63, 11.1, 0.38, GREEN, 14, bold=True, italic=True, align=PP_ALIGN.CENTER)
    labels = [
        ('spike(t)', '{0,1}\nδ-mod', BLUE, BLCK),
        ('FC1\nLinear', 'W₁:[32,1]\nno bias', BLUE, DBLU),
        ('LIF1\nLeaky', 'β learned\nmem persists', GREEN, DGRN),
        ('FC2\nLinear', 'W₂:[2,32]\nno bias', BLUE, DBLU),
        ('Σ logit', '+= FC2(spk₁)\n200 steps', ORA, PANEL),
        ('argmax', 'Normal\nFault', GREEN, DGRN),
    ]
    x0, y0, bw, bh, gap = 0.3, 2.25, 1.85, 1.8, 0.3
    for i,(lbl,sub,lc,fc) in enumerate(labels):
        bx = x0+i*(bw+gap)
        R(s, bx, y0, bw, bh, fc, lc, Pt(2.5))
        T(s, lbl, bx+0.1, y0+0.25, bw-0.2, 0.65, TXT, 14, bold=True, align=PP_ALIGN.CENTER)
        T(s, sub, bx+0.1, y0+0.98, bw-0.2, 0.65, GRAY, 11, align=PP_ALIGN.CENTER)
        if i<len(labels)-1:
            L(s, bx+bw+0.02, y0+bh/2, bx+bw+gap-0.04, y0+bh/2, BLUE, Pt(2))
    # Code pseudocode
    R(s, 0.3, 4.25, 12.73, 1.85, BLCK, BORD)
    T(s, '// Per timestep (5 ms, 200 Hz on ESP32):', 0.5, 4.35, 12.3, 0.32, DGR, 11)
    T(s, 'spike = (|az−prev| > θ) ? 1.0 : 0.0', 0.5, 4.65, 6.0, 0.35, GREEN, 12, italic=True)
    T(s, 'for j in 0..31:', 0.5, 4.98, 6.0, 0.32, BLUE, 12, italic=True)
    T(s, '  mem1[j] = β·mem1[j] + spike·W1[j]', 0.5, 5.28, 6.0, 0.32, BLUE, 12, italic=True)
    T(s, '  if mem1[j]≥θ: { mem1[j]−=θ; out+=W2[j]; }  // no matrix multiply for binary input',
      0.5, 5.58, 12.0, 0.35, GREEN, 12, italic=True)
    T(s, 'if dec_count ≥ 200: classify = argmax(out_acc)', 6.5, 4.65, 6.3, 0.35, ORA, 12, italic=True)
    R(s, 3.5, 6.25, 6.33, 0.6, DGRN, GREEN, Pt(2))
    T(s, 'W₁(32) + W₂(64) + β(1) = 97 params total', 3.7, 6.37, 5.9, 0.38, GREEN, 16, bold=True, align=PP_ALIGN.CENTER)

def ann_s7(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Surrogate Gradient', 'ANN Course · SNN Predictive Maintenance')
    R(s, 0.4, 1.55, 12.53, 0.58, DRED, RED, Pt(1.5))
    T(s, 'Problem: ∂H(mem−θ)/∂mem = 0 everywhere  →  zero gradient  →  no learning!',
      0.55, 1.65, 12.2, 0.38, RED, 15, bold=True, align=PP_ALIGN.CENTER)
    # Left: Forward (step function)
    R(s, 0.3, 2.25, 5.9, 4.55, BLCK, BORD)
    T(s, 'Forward Pass', 0.5, 2.35, 5.5, 0.38, GRAY, 14, bold=True, align=PP_ALIGN.CENTER)
    T(s, 'spike = H(mem − θ)  (Heaviside)', 0.5, 2.72, 5.5, 0.35, ORA, 13, italic=True, align=PP_ALIGN.CENTER)
    # Step function visualization (manual)
    L(s, 0.7, 5.8, 6.0, 5.8, BORD, Pt(1))
    L(s, 3.35, 2.9, 3.35, 6.1, BORD, Pt(1))
    L(s, 0.7, 5.8, 3.35, 5.8, ORA, Pt(3))   # y=0 part
    L(s, 3.35, 3.2, 6.0, 3.2, ORA, Pt(3))   # y=1 part
    L(s, 3.35, 3.2, 3.35, 5.8, ORA, Pt(3))  # vertical jump
    T(s, 'H(x)', 3.5, 2.95, 1.0, 0.3, ORA, 12)
    T(s, '1', 0.45, 3.15, 0.3, 0.3, DGR, 12)
    T(s, '0', 0.45, 5.72, 0.3, 0.3, DGR, 12)
    T(s, 'grad=0', 1.2, 5.3, 1.5, 0.3, RED, 11)
    T(s, 'grad=0', 4.2, 3.0, 1.5, 0.3, RED, 11)
    T(s, 'mem−θ', 5.7, 5.85, 0.85, 0.3, DGR, 11)
    # Right: Backward (surrogate bell)
    R(s, 6.7, 2.25, 6.2, 4.55, BLCK, BORD)
    T(s, 'Backward Pass  (Surrogate)', 6.9, 2.35, 5.8, 0.38, GRAY, 14, bold=True, align=PP_ALIGN.CENTER)
    T(s, "grad ≈ σ'(slope·(mem−θ))", 6.9, 2.72, 5.8, 0.35, GREEN, 13, italic=True, align=PP_ALIGN.CENTER)
    L(s, 6.9, 5.8, 12.7, 5.8, BORD, Pt(1))
    L(s, 9.8, 2.9, 9.8, 6.1, BORD, Pt(1))
    # Bell curve (triangle approximation with line segments)
    pts = [(6.9,5.78),(7.6,5.6),(8.2,5.1),(8.9,4.2),(9.3,3.35),(9.8,2.98),
           (10.3,3.35),(10.7,4.2),(11.4,5.1),(12.0,5.6),(12.7,5.78)]
    for i in range(len(pts)-1):
        L(s, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], GREEN, Pt(2.5))
    T(s, 'peak', 9.6, 2.72, 1.0, 0.3, GREEN, 11)
    T(s, 'smooth everywhere', 10.1, 4.2, 2.5, 0.35, GREEN, 11)
    T(s, 'mem−θ', 12.4, 5.85, 0.9, 0.3, DGR, 11)
    T(s, '→ replace gradient', 6.1, 4.2, 0.85, 0.5, GRAY, 10, align=PP_ALIGN.CENTER)
    # Formula
    R(s, 0.4, 6.65, 12.53, 0.62, PANEL, BORD)
    T(s, "σ'(x) = slope / (2·(1 + |slope·x|)²)    slope=5  →  peak=1.25 at x=0",
      0.6, 6.77, 12.2, 0.38, GREEN, 15, italic=True, align=PP_ALIGN.CENTER)

def ann_s8(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Truncated BPTT  (TBPTT)', 'ANN Course · SNN Predictive Maintenance')
    R(s, 0.3, 1.55, 6.2, 0.58, DRED, RED, Pt(1))
    T(s, 'Full BPTT: O(T) memory, T=120,000 → impractical', 0.5, 1.65, 5.9, 0.38, RED, 13, bold=True)
    R(s, 6.8, 1.55, 6.2, 0.58, DGRN, GREEN, Pt(1.5))
    T(s, 'TBPTT: unroll C=256 steps · detach at boundary', 7.0, 1.65, 5.9, 0.38, GREEN, 13, bold=True)
    # Timeline
    T(s, 'Training: stream divided into chunks (T=256 each)', 0.5, 2.25, 12.33, 0.35, GRAY, 13, bold=True, align=PP_ALIGN.CENTER)
    colors = [BLUE,GREEN,BLUE,GREEN,BLUE]
    for i,c in enumerate(colors):
        R(s, 0.5+i*2.4, 2.65, 2.1, 0.75, DBLU if c==BLUE else DGRN, c, Pt(2))
        T(s, f'chunk {i+1}\nT=256', 0.55+i*2.4, 2.72, 2.0, 0.6, c, 12, bold=True, align=PP_ALIGN.CENTER)
    R(s, 12.2, 2.65, 0.9, 0.75, PANEL, BORD)
    T(s, '...', 12.25, 2.85, 0.8, 0.35, DGR, 18, align=PP_ALIGN.CENTER)
    # Gradient arrows (within chunk)
    for i in range(4):
        L(s, 1.55+i*2.4, 2.65, 1.55+i*2.4, 2.3, GREEN, Pt(1.5))
    T(s, '∂L/∂W', 0.55, 2.12, 2.0, 0.3, GREEN, 11, align=PP_ALIGN.CENTER)
    T(s, '∂L/∂W', 2.95, 2.12, 2.0, 0.3, GREEN, 11, align=PP_ALIGN.CENTER)
    # Detach marks
    for i in range(4):
        T(s, '✂', 2.45+i*2.4, 2.9, 0.4, 0.35, RED, 14, align=PP_ALIGN.CENTER)
    # mem persists
    R(s, 0.4, 3.55, 11.8, 0.45, DBLU, BLUE, Pt(1), st=RECT)
    T(s, 'mem₁ persistent  →  carries temporal context across ALL chunks during INFERENCE',
      0.5, 3.6, 11.7, 0.35, BLUE, 13, bold=True, align=PP_ALIGN.CENTER)
    # Code
    R(s, 0.3, 4.15, 6.8, 2.0, BLCK, BORD)
    T(s, 'Training code:', 0.5, 4.25, 6.5, 0.3, DGR, 11, bold=True)
    for i,(line_,c) in enumerate([
        ('for chunk in chunks:', GREEN),
        ('  mem1 = mem1.detach()   # stop gradient', GRAY),
        ('  out, mem1 = model(chunk, mem1)', BLUE),
        ('  loss = criterion(out, label)', ORA),
        ('  loss.backward(); optimizer.step()', GREEN),
    ]):
        T(s, line_, 0.5, 4.55+i*0.32, 6.5, 0.3, c, 12, italic=True)
    # Trade-offs
    R(s, 7.4, 4.15, 5.6, 2.0, PANEL, BORD)
    T(s, 'Trade-offs', 7.6, 4.25, 5.2, 0.35, ORA, 13, bold=True)
    for i,(icon,text,c) in enumerate([
        ('✓','Memory: O(256) not O(120,000)', GREEN),
        ('✓','Stable on long sequences', GREEN),
        ('✗','Loses dependencies > 256 steps', RED),
        ('→','mem₁ bridges long context at inference', BLUE),
    ]):
        T(s, icon, 7.6, 4.62+i*0.38, 0.3, 0.35, c, 12, bold=True)
        T(s, text, 7.95, 4.62+i*0.38, 4.8, 0.35, GRAY, 12)

def ann_s9(prs):
    ml_s9(prs)  # same results slide

def ann_s10(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Edge Deployment — ESP32', 'ANN Course · SNN Predictive Maintenance')
    # Board diagram (left)
    R(s, 0.3, 1.55, 3.5, 5.5, RGBColor(0x0a,0x1a,0x0a), GREEN, Pt(2))
    T(s, 'ESP32 NodeMCU-32S', 0.5, 1.65, 3.1, 0.38, GREEN, 12, bold=True, align=PP_ALIGN.CENTER)
    R(s, 0.8, 2.15, 2.5, 1.4, PANEL, BORD)
    T(s, 'Xtensa LX6\n240 MHz · 520 KB RAM\nFlash: 4 MB', 0.95, 2.28, 2.2, 1.0, GRAY, 11, align=PP_ALIGN.CENTER)
    R(s, 0.6, 3.75, 1.4, 0.8, RGBColor(0x1a,0x1a,0x2a), BLUE, Pt(1.5))
    T(s, 'MPU6050\nI2C·200Hz', 0.7, 3.88, 1.2, 0.55, BLUE, 10, align=PP_ALIGN.CENTER)
    R(s, 2.2, 3.75, 1.4, 0.8, RGBColor(0x1a,0x1a,0x0a), ORA, Pt(1.5))
    T(s, 'G · R · B\nLEDs', 2.3, 3.88, 1.2, 0.55, ORA, 10, align=PP_ALIGN.CENTER)
    R(s, 0.8, 4.75, 2.5, 0.6, PANEL, BORD)
    T(s, 'Serial 115200 → Dashboard', 0.95, 4.88, 2.2, 0.35, GRAY, 10, align=PP_ALIGN.CENTER)
    T(s, '● Green = Normal\n● Red flash = Fault', 0.5, 5.55, 3.1, 0.6, GRAY, 11)
    # Code (right)
    R(s, 4.1, 1.55, 8.9, 5.75, BLCK, BORD)
    code_lines = [
        ('// firmware/src/main.cpp  (-DINFERENCE_MODE)', DGR),
        ('static float  mem1[SNN_HIDDEN] = {};  // never reset', GRAY),
        ('static float  out_acc[2]       = {};  // fault/normal', GRAY),
        ('', GRAY),
        ('void loop() {', DGR),
        ('  float az = mpu.getAccelZ();', GRAY),
        ('  // ─ delta encode ─', BLUE),
        ('  float spike = (|az-prev|>THRESH) ? 1.0f : 0.0f;', ORA),
        ('  // ─ FC1 + LIF1 (mem never zeroed) ─', BLUE),
        ('  for (int j=0; j<SNN_HIDDEN; j++) {', GREEN),
        ('    mem1[j] = BETA*mem1[j] + spike*W1[j];', GREEN),
        ('    if (mem1[j]>=THRESH1) {', GREEN),
        ('      mem1[j] -= THRESH1;         // soft reset', GREEN),
        ('      out_acc[0] += W2[0*H+j];   // fault', RED),
        ('      out_acc[1] += W2[1*H+j];   // normal', BLUE),
        ('    }', GREEN),
        ('  }', GREEN),
        ('  // ─ classify every 200 samples ─', ORA),
        ('  if (++dec >= INTERVAL) { classify(); reset(); }', ORA),
        ('}', DGR),
        ('// Model: 0.38 KB · binary input = no matrix multiply', DGR),
    ]
    for i,(line_,c) in enumerate(code_lines):
        T(s, line_, 4.25, 1.65+i*0.265, 8.65, 0.27, c, 10, italic=True)

def ann_s11(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'ANN vs SNN — Full Comparison', 'ANN Course · SNN Predictive Maintenance')
    tbl = s.shapes.add_table(8, 3, Inches(0.4), Inches(1.58), Inches(12.53), Inches(5.7)).table
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(4.85)
    tbl.columns[2].width = Inches(4.88)
    col_headers = ['Property', 'Traditional ANN (MLP)', 'Streaming SNN (this work)']
    col_colors  = [BORD, ORA, GREEN]
    row_data = [
        ('Activation',        'float32  ∈ [0, 1]',           'binary  {0, 1}  (spike)',       GREEN),
        ('FC1 input size',    'WINDOW = 128  ← fixed!',       '1  (one spike per timestep)',   GREEN),
        ('Parameters',        '128×32 + 32×2 = 4,160',        '32 + 64 + 1 = 97',             GREEN),
        ('Temporal Memory',   'none  (per-window, stateless)', 'mem₁ persistent across all t', GREEN),
        ('Training',          'Standard BPTT',                 'Surrogate grad + TBPTT',        BLUE),
        ('Real-time capable', 'windowed only  (batch)',         'true streaming  (per-sample)',  GREEN),
        ('Test Accuracy',     '~99%  (4160 params)',           '100%  (97 params)',             GREEN),
    ]
    for ci,(h,hc) in enumerate(zip(col_headers,col_colors)):
        cell = tbl.cell(0,ci); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold=True; p.runs[0].font.color.rgb=hc
        p.runs[0].font.size=Pt(13); p.alignment=PP_ALIGN.CENTER
    for ri,(prop,ann_v,snn_v,vc) in enumerate(row_data):
        fc = BLCK if ri%2==0 else PANEL
        for ci,(val,color) in enumerate([(prop,GRAY),(ann_v,TXT),(snn_v,vc)]):
            cell = tbl.cell(ri+1,ci); cell.text=val
            cell.fill.solid(); cell.fill.fore_color.rgb = (DGRN if ri==6 and ci==2 else fc)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.color.rgb=color; p.runs[0].font.size=Pt(12)
            p.runs[0].font.bold=(ci==2); p.alignment=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT

def ann_s12(prs):
    s = blank(prs); set_bg(s)
    hdr(s, 'Conclusion', 'ANN Course · SNN Predictive Maintenance')
    T(s, 'SNN Contributions', 0.5, 1.55, 5.8, 0.38, GREEN, 15, bold=True)
    for i,t in enumerate([
        'LIF neuron: built-in temporal integration',
        'Surrogate gradient enables BPTT on spikes',
        'Streaming: one weight per hidden unit',
        'TBPTT: scalable training on long streams',
        'Edge deployment: 97 params in C on ESP32',
        '100% test accuracy on CWRU benchmark',
    ]):
        check_item(s, t, 0.5, 2.0+i*0.7, GREEN)
    L(s, 6.5, 1.55, 6.5, 7.1, BORD, Pt(1.5))
    T(s, 'Key Insights', 6.8, 1.55, 6.1, 0.38, BLUE, 15, bold=True)
    insights = [
        ('Neuromorphic ≠ just biological', 'sparse spikes → real computational benefit'),
        ('Delta mod as feature engineering', 'maps domain knowledge to spike rates'),
        ('Persistent state = temporal context', 'replaces explicit window or LSTM cell'),
        ('Soft reset preserves energy', 'more info than hard reset to zero'),
    ]
    for i,(t,d) in enumerate(insights):
        bullet(s, t, 6.8, 2.0+i*1.12)
        T(s, d, 7.15, 2.32+i*1.12, 5.8, 0.3, DGR, 13)
    R(s, 0.4, 6.65, 12.53, 0.7, DGRN, GREEN, Pt(2.5))
    T(s, 'Streaming SNN: neuromorphic principles yield practical efficiency gains  —  97 params, 100% accuracy',
      0.6, 6.77, 12.2, 0.42, GREEN, 14, bold=True, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  BUILD
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def build_ml():
    prs = new_prs()
    for fn in [ml_s1,ml_s2,ml_s3,ml_s4,ml_s5,ml_s6,
               ml_s7,ml_s8,ml_s9,ml_s10,ml_s11,ml_s12]:
        fn(prs)
    out = OUT/'ml_presentation.pptx'
    prs.save(str(out))
    print(f'OK  {out}')

def build_ann():
    prs = new_prs()
    for fn in [ann_s1,ann_s2,ann_s3,ann_s4,ann_s5,ann_s6,
               ann_s7,ann_s8,ann_s9,ann_s10,ann_s11,ann_s12]:
        fn(prs)
    out = OUT/'ann_presentation.pptx'
    prs.save(str(out))
    print(f'OK  {out}')

if __name__ == '__main__':
    build_ml()
    build_ann()
    print('Done.')

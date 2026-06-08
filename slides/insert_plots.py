"""
Insert CV result plots into the ML and ANN PPTX presentations.
Replaces the Results slide with actual CV plots.
Run from slides/ directory:  python insert_plots.py
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PLOTS   = pathlib.Path("../training/plots")
ML_IN   = pathlib.Path("ml_presentation.pptx")
ANN_IN  = pathlib.Path("ann_presentation.pptx")

BG    = RGBColor(0x0d, 0x11, 0x17)
TXT   = RGBColor(0xf0, 0xf6, 0xfc)
GREEN = RGBColor(0x3f, 0xb9, 0x50)
GRAY  = RGBColor(0x8b, 0x94, 0x9e)
BORD  = RGBColor(0x30, 0x36, 0x3d)
DGR   = RGBColor(0x48, 0x4f, 0x58)
BLUE  = RGBColor(0x58, 0xa6, 0xff)


def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG


def T(slide, text, l, t, w, h, color=TXT, size=14, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.color.rgb = color; r.font.size = Pt(size); r.font.bold = bold


def add_results_slide(prs, idx):
    """Replace slide at idx with a new CV results slide."""
    # Remove the old slide
    xml_slides = prs.slides._sldIdLst
    old_slide_ref = xml_slides[idx]
    xml_slides.remove(old_slide_ref)

    # Add new blank slide at the end, then reorder
    layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(layout)
    set_bg(new_slide)

    # Move to position idx
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
    prs.slides._sldIdLst.insert(idx, old_slide_ref.__class__())
    # Simpler: just rebuild from scratch and append at end (accept different position)
    return new_slide


def build_cv_results_slide(prs):
    """Add a CV results slide at the end of the presentation."""
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_bg(slide)

    # Header
    T(slide, 'Machine Learning · SNN Predictive Maintenance',
      0.3, 0.18, 9, 0.32, DGR, 11)
    T(slide, '5-Fold Cross-Validation Results',
      0.5, 0.52, 12.33, 0.85, TXT, 34, bold=True, align=PP_ALIGN.CENTER)
    # header line
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                   Inches(0), Inches(1.45), Inches(13.33), Inches(1.45))
    c.line.color.rgb = BORD; c.line.width = Pt(1)

    # CV Summary table
    T(slide, 'Summary (mean ± std across 5 folds)', 0.4, 1.55, 5.5, 0.38, GRAY, 13, bold=True)
    rows = [
        ('Accuracy',   '1.0000 ± 0.0000', GREEN),
        ('Precision',  '1.0000 ± 0.0000', GREEN),
        ('Recall',     '1.0000 ± 0.0000', GREEN),
        ('F1 Score',   '1.0000 ± 0.0000', GREEN),
        ('ROC-AUC',    '1.0000 ± 0.0000', GREEN),
    ]
    for i, (name, val, c_) in enumerate(rows):
        from pptx.util import Inches as In
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        sh = slide.shapes.add_shape(1, In(0.4), In(1.95+i*0.58), In(5.5), In(0.5))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x0f,0x2a,0x1a) if i%2==0 else RGBColor(0x0a,0x14,0x0a)
        sh.line.color.rgb = BORD; sh.line.width = Pt(1)
        T(slide, name, 0.55, 2.03+i*0.58, 2.2, 0.38, GRAY,  13)
        T(slide, val,  2.8,  2.03+i*0.58, 3.0, 0.38, c_,    14, bold=True)

    # Classification report summary
    T(slide, 'All 1425 chunks  ·  473 fault  ·  952 normal',
      0.4, 5.0, 5.5, 0.35, DGR, 12)
    T(slide, 'fault   P=1.00  R=1.00  F1=1.00  (473 chunks)',
      0.4, 5.38, 5.5, 0.35, GREEN, 12, bold=True)
    T(slide, 'normal  P=1.00  R=1.00  F1=1.00  (952 chunks)',
      0.4, 5.73, 5.5, 0.35, BLUE,  12, bold=True)

    # Insert plots
    conv_p = PLOTS / 'convergence.png'
    cm_p   = PLOTS / 'confusion_matrix.png'
    roc_p  = PLOTS / 'roc_curve.png'

    if conv_p.exists():
        slide.shapes.add_picture(str(conv_p), Inches(6.0), Inches(1.55), Inches(7.0), Inches(3.5))
        T(slide, 'Training Convergence', 6.0, 5.1, 7.0, 0.3, GRAY, 11, align=PP_ALIGN.CENTER)

    if cm_p.exists():
        slide.shapes.add_picture(str(cm_p), Inches(6.0), Inches(5.45), Inches(3.5), Inches(1.8))

    if roc_p.exists():
        slide.shapes.add_picture(str(roc_p), Inches(9.7), Inches(5.45), Inches(3.4), Inches(1.8))

    return slide


def update_pptx(path_in, path_out):
    prs = Presentation(str(path_in))
    build_cv_results_slide(prs)
    prs.save(str(path_out))
    print(f"Saved: {path_out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    import sys
    sys.path.insert(0, str(pathlib.Path(__file__).parent.parent / "slides"))

    update_pptx(ML_IN,  ML_IN.with_name("ml_presentation_cv.pptx"))
    update_pptx(ANN_IN, ANN_IN.with_name("ann_presentation_cv.pptx"))
    print("Done.")
